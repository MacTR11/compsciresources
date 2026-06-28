#!/usr/bin/env python3
"""
Build editable PowerPoint lesson decks (.pptx) for every subtopic from the
deck-data JSON in source/teacher-toolkit/deck-data/.

Each deck is a ready-to-teach skeleton: title, objectives, a "Do Now"
retrieval starter, keywords, content slides, a worked example, a task, a
challenge and an exit ticket (with answers) — all editable in PowerPoint /
Google Slides / Keynote.

Usage:  python3 build/build-powerpoints.py
"""
import sys, os, glob, json

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
except ImportError:
    sys.exit("python-pptx not installed. Run: pip install python-pptx")

REPO = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
DATA = os.path.join(REPO, "source", "teacher-toolkit", "deck-data")
OUT = os.path.join(REPO, "Teacher-Toolkit", "Lesson-PowerPoints")

NAVY = RGBColor(0x1A, 0x36, 0x5D)
BLUE = RGBColor(0x2B, 0x6C, 0xB0)
DARK = RGBColor(0x22, 0x22, 0x22)
GREY = RGBColor(0x66, 0x66, 0x66)
AMBER = RGBColor(0xB7, 0x6E, 0x00)
W, H = Inches(13.333), Inches(7.5)


def _bg(slide, rgb):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb


def _box(slide, l, t, w, h):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    return tf


def _set(p, text, size, color, bold=False, italic=False):
    p.text = text
    for r in p.runs:
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.italic = italic
        r.font.name = "Calibri"


def _bullet(tf, text, size=20, color=DARK, level=0, bold=False):
    p = tf.add_paragraph() if tf.paragraphs[0].text or len(tf.paragraphs) > 1 or tf.paragraphs[0].runs else tf.paragraphs[0]
    p.level = level
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold; r.font.name = "Calibri"
    return p


def title_slide(prs, deck):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, NAVY)
    tf = _box(s, Inches(0.8), Inches(2.4), Inches(11.7), Inches(2.0))
    _set(tf.paragraphs[0], deck["subtopic"], 40, RGBColor(0xFF, 0xFF, 0xFF), bold=True)
    p = tf.add_paragraph(); _set(p, deck.get("spec", "OCR H446"), 20, RGBColor(0xBE, 0xD3, 0xEC))
    foot = _box(s, Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.6))
    _set(foot.paragraphs[0], "Lesson deck — editable skeleton: adapt freely.", 14, RGBColor(0x9A, 0xB4, 0xD6), italic=True)


def header(slide, text, color=BLUE):
    bar = _box(slide, Inches(0.6), Inches(0.3), Inches(12.1), Inches(0.9))
    _set(bar.paragraphs[0], text, 30, color, bold=True)


def content_slide(prs, title, fill_fn):
    s = prs.slides.add_slide(prs.slide_layouts[6]); _bg(s, RGBColor(0xFF, 0xFF, 0xFF))
    header(s, title)
    tf = _box(s, Inches(0.7), Inches(1.4), Inches(11.9), Inches(5.6))
    fill_fn(tf)
    return s


def qa_slide(prs, title, pairs, show_answers=True, qsize=22):
    def fill(tf):
        first = True
        for i, qa in enumerate(pairs, 1):
            p = tf.paragraphs[0] if first else tf.add_paragraph(); first = False
            r = p.add_run(); r.text = f"{i}. {qa['q']}"
            r.font.size = Pt(qsize); r.font.bold = True; r.font.color.rgb = DARK; r.font.name = "Calibri"
            if show_answers:
                a = tf.add_paragraph(); a.level = 1
                ra = a.add_run(); ra.text = "Answer: " + qa["a"]
                ra.font.size = Pt(18); ra.font.color.rgb = AMBER; ra.font.italic = True; ra.font.name = "Calibri"
    content_slide(prs, title, fill)


def build_deck(deck, path):
    prs = Presentation(); prs.slide_width = W; prs.slide_height = H
    title_slide(prs, deck)

    # Objectives
    def obj(tf):
        _bullet(tf, "By the end of this lesson you can:", 22, BLUE, bold=True)
        for o in deck.get("objectives", []):
            _bullet(tf, o, 22, DARK, level=1)
    content_slide(prs, "Learning objectives", obj)

    # Do Now retrieval (questions then answers slide)
    if deck.get("retrieval"):
        qa_slide(prs, "Do Now — retrieval starter", deck["retrieval"], show_answers=False)
        qa_slide(prs, "Do Now — answers", deck["retrieval"], show_answers=True)

    # Keywords
    if deck.get("keywords"):
        def kw(tf):
            for k in deck["keywords"]:
                p = tf.paragraphs[0] if (tf.paragraphs[0].runs == [] and len(tf.paragraphs) == 1 and not tf.paragraphs[0].text) else tf.add_paragraph()
                r = p.add_run(); r.text = k["term"]; r.font.bold = True; r.font.size = Pt(20); r.font.color.rgb = NAVY; r.font.name = "Calibri"
                r2 = p.add_run(); r2.text = " — " + k["def"]; r2.font.size = Pt(18); r2.font.color.rgb = DARK; r2.font.name = "Calibri"
        content_slide(prs, "Key terminology", kw)

    # Content chunks
    for chunk in deck.get("content", []):
        def fill(tf, chunk=chunk):
            for pt in chunk.get("points", []):
                _bullet(tf, pt, 22, DARK)
        content_slide(prs, chunk.get("heading", "Content"), fill)

    # Worked example
    we = deck.get("worked_example")
    if we and we.get("steps"):
        def fill(tf, we=we):
            _bullet(tf, we.get("title", "Worked example"), 22, BLUE, bold=True)
            for i, st in enumerate(we["steps"], 1):
                _bullet(tf, f"{i}. {st}", 20, DARK, level=1)
        content_slide(prs, "Worked example", fill)

    # Task + Challenge
    def task(tf):
        _bullet(tf, "Task", 24, BLUE, bold=True)
        _bullet(tf, deck.get("task", ""), 22, DARK, level=1)
        if deck.get("challenge"):
            c = tf.add_paragraph(); rc = c.add_run(); rc.text = "Challenge / stretch"
            rc.font.size = Pt(24); rc.font.bold = True; rc.font.color.rgb = AMBER; rc.font.name = "Calibri"
            _bullet(tf, deck["challenge"], 22, DARK, level=1)
    content_slide(prs, "Your turn", task)

    # Exit ticket
    if deck.get("exit_ticket"):
        qa_slide(prs, "Exit ticket", deck["exit_ticket"], show_answers=False)
        qa_slide(prs, "Exit ticket — answers", deck["exit_ticket"], show_answers=True)

    os.makedirs(os.path.dirname(path), exist_ok=True)
    prs.save(path)


def main():
    files = sorted(glob.glob(os.path.join(DATA, "*.json")))
    if not files:
        sys.exit("No deck-data JSON found — run the deck-content writers first.")
    n = 0
    for jf in files:
        deck = json.load(open(jf, encoding="utf-8"))
        out = os.path.join(OUT, os.path.splitext(os.path.basename(jf))[0] + ".pptx")
        build_deck(deck, out)
        print(f"  {os.path.basename(out)}")
        n += 1
    print(f"\nDone. {n} lesson decks in {os.path.relpath(OUT, REPO)}")


if __name__ == "__main__":
    main()
